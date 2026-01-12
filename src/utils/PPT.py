from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# ---------- CONFIG -------------
prs = Presentation()
prs.slide_width = Inches(13.33)  # 16:9 ratio
prs.slide_height = Inches(7.5)

# ---------- SLIDE CONTENT ----------
slides = [
    ("Datakwaliteit & Transformatie met dbt",
     ["Van ruwe geodata naar betrouwbare inzichten voor beleid, vastgoed en GIS",
      "Digitale Kadastrale Kaart – end-to-end ETL pipeline"],
     "Welkom. Ik ben Paschal Bongyu. Vandaag presenteer ik hoe we ruwe geodata omzetten naar betrouwbare inzichten."),

    ("Waarom Datakwaliteit Telt",
     ["Goede data → betere beslissingen",
      "Slechte data → fouten, risico’s & compliance-issues",
      "dbt = geautomatiseerde kwaliteitscontrole"],
     "Datakwaliteit is cruciaal. dbt borgt kwaliteit direct aan de bron."),

    ("Projectoverzicht",
     ["Doel: bedrijfslogica + kwaliteitscontrole",
      "Aanpak: modulaire dbt-transformaties",
      "Tools: Spark, Delta Lake, dbt"],
     "Een schaalbare, transparante en reproduceerbare aanpak."),

    ("Bronze-laag ingestie",
     ["GML → Spark ingestiescript",
      "Delta table: dpt_test_bronze.bron.kadastraal_data_delta",
      "Centrale opslag (‘single source of truth’)"],
     "De Bronze-laag verzamelt alle ruwe data samen."),

    ("dbt Configuratie & Authenticatie",
     ["dbt init → PAT failed",
      "Azure CLI login → az login",
      "profiles.yml geconfigureerd",
      "Verbinding geslaagd"],
     "We overwonnen PAT-beperkingen met Azure OAuth."),

    ("Tool Setup",
     ["Python → python --version",
      "dbt-databricks → pip install dbt-databricks",
      "Git → git --version"],
     "Essentiële tools lokaal geïnstalleerd en getest."),

    ("Transformatie Modellen",
     ["Flattening van GML-structuren",
      "Aggregatie per gemeente",
      "Filter: oppervlakte > 0"],
     "We zetten complexe geodata om in analyseerbare structuren."),

    ("Validatiestrategie met dbt Tests",
     ["not_null, unique, accepted_values",
      "Custom: test_oppervlakte_positive",
      "11 van 11 tests geslaagd"],
     "Onze data voldoet aan alle kwaliteitsregels."),

    ("Resultaten & Impact",
     ["Succesvolle modelbouw",
      "Geen duplicaten of nulls",
      "Data klaar voor BI & GIS analyse"],
     "Een robuuste en betrouwbare datapipeline."),

    ("Pipeline Overzicht",
     ["GML → Spark → Delta",
      "dbt Models → dbt Tests → BI Tools",
      "Volledig reproduceerbaar"],
     "Elke stap is gecontroleerd, getest en transparant."),

    ("Business Impact",
     ["Betrouwbare inzichten voor beleid",
      "Minder risico’s",
      "Snellere time-to-insight"],
     "De pipeline levert strategische waarde voor Kadaster."),

    ("Vragen & Afsluiting",
     ["Bedankt voor uw aandacht",
      "Vragen zijn welkom"],
     "Dank u wel. Ik beantwoord graag uw vragen.")
]

# ---------- SLIDE GENERATOR ----------
for title, bullets, notes in slides:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # lege layout

    # Titel
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(11), Inches(1.5))
    tf = tb.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(40)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # zwart

    # Bulletpoints
    left = Inches(0.9)
    top = Inches(2.0)
    width = Inches(11)
    height = Inches(4)
    tx_box = slide.shapes.add_textbox(left, top, width, height)
    tf2 = tx_box.text_frame
    tf2.word_wrap = True
    for b in bullets:
        p = tf2.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(26)
        p.font.color.rgb = RGBColor(0, 0, 0)  # zwart

    # Spreker-notities
    slide.notes_slide.notes_text_frame.text = notes

# ---------- SAVE -----------
output_path = "Kadaster_DBT_Pipeline_Presentatie.pptx"
prs.save(output_path)

print(f"PowerPoint succesvol aangemaakt: {output_path}")
